#!/usr/bin/env python3
"""Generate the 'Value of Prescribed Data' conference talk deck (.pptx)."""
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY   = RGBColor(0x12, 0x23, 0x3F)   # primary dark
INK    = RGBColor(0x22, 0x2A, 0x33)   # body text
TEAL   = RGBColor(0x14, 0x8F, 0x8F)   # accent / "what matters"
AMBER  = RGBColor(0xC0, 0x55, 0x2B)   # "what doesn't matter" / caution
MUTED  = RGBColor(0x5A, 0x60, 0x6A)   # secondary text
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # panel bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RULE   = RGBColor(0xD5, 0xDB, 0xE1)

FONT = "Calibri"
FONT_H = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- inline text parser: **bold**, ==accent bold==, //muted// ----------
def _segments(text):
    # split preserving markers
    pattern = re.compile(r'(\*\*.+?\*\*|==.+?==|//.+?//|\*[^*]+?\*)')
    segs = []
    for part in pattern.split(text):
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            segs.append((part[2:-2], dict(bold=True, color=INK)))
        elif part.startswith('==') and part.endswith('=='):
            segs.append((part[2:-2], dict(bold=True, color=TEAL)))
        elif part.startswith('//') and part.endswith('//'):
            segs.append((part[2:-2], dict(bold=False, color=MUTED, italic=True)))
        elif part.startswith('*') and part.endswith('*') and len(part) > 2:
            # single-asterisk emphasis: italic, inherit paragraph color
            segs.append((part[1:-1], dict(italic=True)))
        else:
            segs.append((part, dict()))
    return segs


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _txbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def _set_para(p, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT,
              space_after=6, space_before=0, italic=False, line=None):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    segs = _segments(text) if ('*' in text or '==' in text or '//' in text) else [(text, {})]
    for seg_text, style in segs:
        r = p.add_run()
        r.text = seg_text
        r.font.size = Pt(size)
        r.font.name = font
        r.font.bold = style.get('bold', bold)
        r.font.italic = style.get('italic', italic)
        r.font.color.rgb = style.get('color', color)
    return p


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _page_num(slide, n):
    box, tf = _txbox(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    _set_para(tf.paragraphs[0], str(n), 11, MUTED, align=PP_ALIGN.RIGHT, space_after=0)


def _kicker(slide, text):
    box, tf = _txbox(slide, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set_para(tf.paragraphs[0], text.upper(), 12, TEAL, bold=True, space_after=0)


_page = {"n": 0}


# ---------- slide builders ----------
def title_slide(title, subtitle, presenter, meta):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(bg, NAVY)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.55), Inches(2.2), Inches(0.10))
    _fill(bar, TEAL)
    box, tf = _txbox(slide, Inches(0.9), Inches(2.8), Inches(11.5), Inches(2.6))
    _set_para(tf.paragraphs[0], title, 44, WHITE, bold=True, font=FONT_H, space_after=10, line=1.05)
    _set_para(tf.add_paragraph(), subtitle, 24, RGBColor(0xBF, 0xD8, 0xD8), space_after=0, line=1.05)
    box2, tf2 = _txbox(slide, Inches(0.9), Inches(6.1), Inches(11.5), Inches(1.0))
    _set_para(tf2.paragraphs[0], presenter, 16, WHITE, bold=True, space_after=2)
    _set_para(tf2.add_paragraph(), meta, 13, RGBColor(0x9A, 0xB0, 0xB8), space_after=0)
    return slide


def section_slide(num, title, blurb=""):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(bg, NAVY)
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), EMU_H)
    _fill(side, TEAL)
    box, tf = _txbox(slide, Inches(1.1), Inches(2.7), Inches(11), Inches(2.4))
    _set_para(tf.paragraphs[0], num, 22, TEAL, bold=True, space_after=6)
    _set_para(tf.add_paragraph(), title, 40, WHITE, bold=True, font=FONT_H, space_after=12, line=1.03)
    if blurb:
        _set_para(tf.add_paragraph(), blurb, 18, RGBColor(0xBF, 0xD8, 0xD8), space_after=0, line=1.1)
    return slide


def content_slide(title, kicker=None, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    _page["n"] += 1
    if kicker:
        _kicker(slide, kicker)
    box, tf = _txbox(slide, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set_para(tf.paragraphs[0], title, 30, NAVY, bold=True, font=FONT_H, space_after=2, line=1.0)
    if subtitle:
        _set_para(tf.add_paragraph(), subtitle, 15, MUTED, space_after=0, line=1.05)
    # underline rule
    y = Inches(1.72) if subtitle else Inches(1.55)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), y, Inches(1.5), Pt(3))
    _fill(rule, TEAL)
    _page_num(slide, _page["n"])
    return slide


def bullets(slide, items, top=Inches(2.0), left=Inches(0.75), width=Inches(11.9),
            height=Inches(5.0), base_size=18):
    box, tf = _txbox(slide, left, top, width, height)
    first = True
    for it in items:
        lvl, text = it[0], it[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        size = base_size if lvl == 0 else base_size - 3
        color = INK if lvl == 0 else MUTED
        # bullet glyph
        bullet = "▸ " if lvl == 0 else "– "
        p.level = lvl
        _set_para(p, bullet + text, size, color,
                  space_after=(9 if lvl == 0 else 4),
                  space_before=(5 if (lvl == 0 and not first) else 0),
                  line=1.06)
    return box


def two_col(slide, left_title, left_items, right_title, right_items,
            top=Inches(2.0), lcolor=TEAL, rcolor=AMBER):
    colw = Inches(5.85)
    gap = Inches(0.35)
    lx = Inches(0.75)
    rx = lx + colw + gap
    for x, htxt, items, hc in [(lx, left_title, left_items, lcolor),
                               (rx, right_title, right_items, rcolor)]:
        # header chip
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, colw, Inches(0.5))
        _fill(chip, hc)
        ctf = chip.text_frame
        ctf.margin_top = Pt(2); ctf.margin_bottom = Pt(2)
        _set_para(ctf.paragraphs[0], htxt, 15, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        box, tf = _txbox(slide, x, top + Inches(0.62), colw, Inches(4.4))
        first = True
        for lvl, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            size = 15 if lvl == 0 else 13
            _set_para(p, ("▸ " if lvl == 0 else "   – ") + text, size,
                      INK if lvl == 0 else MUTED, space_after=6, line=1.05)


def result_table(slide, rows, top=Inches(2.05), left=Inches(0.9), width=Inches(11.5),
                 col_ratios=None, highlight_rows=None, height=Inches(4.4)):
    highlight_rows = highlight_rows or {}
    nrows = len(rows)
    ncols = len(rows[0])
    gfx = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = gfx.table
    # column widths
    if col_ratios is None:
        col_ratios = [1.0] * ncols
    tot = sum(col_ratios)
    for i, r in enumerate(col_ratios):
        table.columns[i].width = Emu(int(width * r / tot))
    # kill default styling banding
    tbl = table._tbl
    for tblPr in tbl.iter(qn('a:tblPr')):
        tblPr.set('firstRow', '1')
        tblPr.set('bandRow', '0')
    for ri, row in enumerate(rows):
        table.rows[ri].height = Inches(0.42)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if ri == 0:
                _fill_cell(cell, NAVY)
                _set_para(p, str(val), 13, WHITE, bold=True,
                          align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, space_after=0)
            else:
                hl = highlight_rows.get(ri)
                if hl:
                    _fill_cell(cell, hl[0])
                    txtcolor = hl[1]
                    bold = True
                else:
                    _fill_cell(cell, WHITE if ri % 2 else LIGHT)
                    txtcolor = INK
                    bold = (ci == 0)
                _set_para(p, str(val), 13, txtcolor, bold=bold,
                          align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, space_after=0)
    return table


def _fill_cell(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def note_strip(slide, text, top=Inches(6.75)):
    box, tf = _txbox(slide, Inches(0.75), top, Inches(11.9), Inches(0.55))
    _set_para(tf.paragraphs[0], text, 12.5, MUTED, italic=True, space_after=0, line=1.03)


FIGDIR = "/data/xinyua11/robocasa/talk/figs"

def image_slide(title, kicker, img_path, subtitle=None, bullets_items=None,
                caption=None, img_top=Inches(2.12), max_h=Inches(4.35)):
    """Content slide with a figure (optionally + bullets on the right)."""
    from PIL import Image as _Img
    s = content_slide(title, kicker=kicker, subtitle=subtitle)
    iw, ih = _Img.open(img_path).size
    ar = ih / float(iw)
    if bullets_items:
        w = Inches(7.7)
        h = Emu(int(int(w) * ar))
        if int(h) > int(max_h):
            h = max_h; w = Emu(int(int(h) / ar))
        s.shapes.add_picture(img_path, Inches(0.55), img_top, width=w, height=h)
        bullets(s, bullets_items, top=Inches(2.05), left=Inches(8.45),
                width=Inches(4.35), base_size=15)
    else:
        w = Inches(11.4)
        h = Emu(int(int(w) * ar))
        if int(h) > int(max_h):
            h = max_h; w = Emu(int(int(h) / ar))
        left = Emu(int((int(EMU_W) - int(w)) / 2))
        s.shapes.add_picture(img_path, left, img_top, width=w, height=h)
    if caption:
        note_strip(s, caption)
    return s


# ==========================================================================
# SLIDE 1 — TITLE
# ==========================================================================
title_slide(
    "The Value of Prescribed Data",
    "Insights from my experiments on data-efficient robot imitation learning",
    "Xinyuan",
    "A 20-minute talk  ·  RoboCasa · pi0 / GR00T · 2026",
)
_notes(prs.slides[-1], """
Good [morning/afternoon]. My talk is called 'The Value of Prescribed Data.'

The one-line version: when we train robots and humanoids, the bottleneck is no longer
compute or model size — it's data. And the question I care about is not *how much* data,
but *which* data. 'Prescribed data' means letting the policy itself tell us what to collect next.

I'll give you two things by the end. First, an applied result: prescribed data from a policy's
own failure regions beats random data — but only when you collect it in a concentrated way.
Second, a deeper and more surprising lesson: sophisticated data-selection methods lost to a
trivial heuristic, and I'll explain the exact principle that governs when data selection can work
at all. This is a talk with honest negative results — those turned out to be the most useful part.
""")

# ==========================================================================
# SECTION 1
# ==========================================================================
section_slide("Part 1", "Why prescribed data matters",
              "The bottleneck in training robots is which data, not how much.")
_notes(prs.slides[-1], "Two minutes on motivation: why the *selection* of data is the frontier.")

# Slide: the bottleneck
s = content_slide("The bottleneck is data — and it's expensive",
                  kicker="Part 1 · Motivation")
bullets(s, [
    (0, "Robot & humanoid learning has shifted: model architectures and compute are largely commoditized — ==the scarce resource is demonstration data=="),
    (0, "Every demo has a real cost: teleoperation, human operator time, hardware wear, scene resets"),
    (0, "'Just collect more data' scales cost linearly but returns **diminish** — most new demos repeat what the policy already does well"),
    (1, "A policy at 55% success doesn't need 10,000 more *random* demos — it needs the *right* few hundred"),
    (0, "So the real lever is **data efficiency**: maximize improvement *per added demonstration*"),
])
_notes(prs.slides[-1], """
Start with the pain. In modern robot learning the model is not the hard part anymore — pi0, GR00T,
diffusion policies are all public. What's expensive is data. Every single demonstration costs human
teleoperation time, operator attention, hardware wear, and scene resets.

And here's the trap everyone falls into: 'the policy is at 55%, let's collect more data.' But more
*random* data mostly re-teaches what the policy already knows. The returns diminish fast. If your
policy already grasps mugs reliably, the hundredth mug demo is nearly worthless.

So the objective I care about is data efficiency: improvement per added demonstration. That reframes
the whole problem from 'how much' to 'which.'
""")

# Slide: which data — prescribed loop
s = content_slide("Not how much data — which data",
                  kicker="Part 1 · The idea",
                  subtitle="Prescribed data: the policy tells you what to collect next")
bullets(s, [
    (0, "**Prescribed data** = a closed loop that turns a trained policy into a data-collection spec:"),
    (1, "==Evaluate== the policy → ==find== where it fails → ==prescribe== the data type to add → collect it → retrain → repeat"),
    (0, "Contrast with the status quo: collect a big i.i.d. pool once, train, hope coverage is enough"),
    (0, "Grounded in a real problem: a **Unitree G1 humanoid** learning to **pour** (chemistry pouring) — teleop demos are minutes each, so you cannot collect everything and must choose *which*"),
    (0, "This talk: does prescribing from failure regions actually pay off? And **what makes it work or fail?**"),
], top=Inches(2.1))
_notes(prs.slides[-1], """
Here's the idea in one word: prescription. Instead of collecting a giant i.i.d. dataset once and
hoping, we close a loop. Evaluate the policy. Find where it fails. Prescribe *what kind* of data to
add. Collect exactly that. Retrain. Repeat.

The analogy I like: a doctor doesn't prescribe every drug in the pharmacy — they diagnose, then
prescribe the specific thing. We want the policy to diagnose itself.

This is motivated by a concrete downstream problem — a real Unitree G1 humanoid learning to pour, a
chemistry-pouring task — where each teleoperated demo takes minutes, you genuinely cannot collect
everything, and every demo counts. I'll come back to that real robot near the end.

The two questions for the rest of the talk: (1) does prescribing from failure regions actually beat
just adding random data? and (2) — the more interesting one — what actually determines whether it
works? Because the answer surprised me.
""")

# ==========================================================================
# SECTION 2
# ==========================================================================
section_slide("Part 2", "What people have done so far",
              "Three families: scale it, select it, or correct it — mostly validated off-robot.")
_notes(prs.slides[-1], "Four minutes situating the work. Three families of prior approaches, and the gap each leaves.")

# Slide: scaling
s = content_slide("Approach 1 — scale the data (brute force)",
                  kicker="Part 2 · Prior work")
bullets(s, [
    (0, "**Large cross-embodiment datasets**: Open X-Embodiment / RT-X, DROID, BridgeData — pool demos across labs, tasks, robots"),
    (0, "**Bet**: enough diverse data + a big model ⇒ generalization emerges (the scaling-law playbook from LLMs)"),
    (0, "**Works**, and is the backbone of today's VLAs (pi0, GR00T, OpenVLA are trained this way)"),
    (0, "==But==: cost grows linearly, coverage of the *long tail* stays thin, and it says nothing about **which** demo to collect **next** for a given policy"),
    (1, "Scaling is a strategy for the field; it is not a strategy for *one lab with a fixed budget and a specific weak policy*"),
])
_notes(prs.slides[-1], """
Family one: scale it. Open X-Embodiment, RT-X, DROID, Bridge — the move is to pool enormous amounts
of demonstration data across labs and robots, then train a big model. This is the LLM scaling-law
playbook imported into robotics, and it genuinely works — it's how every VLA you've heard of was built.

But two problems. Cost is linear — double the capability, roughly double the data budget. And the long
tail stays thin no matter how big you go. Most importantly for us: scaling is a field-level strategy.
It gives you no guidance if you're one lab, with a fixed budget, holding one specific policy that fails
in one specific way. It never answers 'which demo next?'
""")

# Slide: selection
s = content_slide("Approach 2 — select the data (be smart about which)",
                  kicker="Part 2 · Prior work")
bullets(s, [
    (0, "**Active learning / uncertainty sampling** — label/collect where the model is least certain (classic, decades deep)"),
    (0, "**Influence functions & gradient attribution** — score each training point by its effect on a target loss:"),
    (1, "LESS, TracIn, DataInf, Datamodels — strong results in **LLM fine-tuning** and **image classification**"),
    (0, "**Coreset / coverage / diversity selection** — a representative or maximally-spread subset"),
    (0, "**Embedding / scene retrieval** — pick demos *visually* similar to failure states (BehaviorRetrieval-style; CLIP / DINOv2 kNN & coverage)"),
    (0, "==The catch==: almost all of this is validated on **held-out loss** in **classification / language** — not on **closed-loop robot success**, and rarely head-to-head"),
])
_notes(prs.slides[-1], """
Family two — the one my work lives in — select it. Be smart about *which* data.

The oldest branch is active learning: collect where the model is uncertain. The branch I lean on
hardest is influence functions and gradient attribution — LESS, TracIn, DataInf, Datamodels. The idea
is elegant: score every candidate training point by how much it would change a target loss, via the
alignment of its gradient with a target gradient. These have real, strong results — in LLM fine-tuning
and image classification. And there's coverage / coreset selection: pick a diverse, representative subset.

Here's the catch that motivated my whole project. Almost all of this evidence is on *held-out loss*,
in *classification or language*. Very little of it has been tested on *closed-loop robot task success*,
where a policy rolls out, compounds its own errors, and either completes the task or doesn't. And these
methods are almost never run head-to-head against each other, or against a dumb baseline. That's the gap.
""")

# Slide: correction + the gap
s = content_slide("Approach 3 — correct the data · and the gap I target",
                  kicker="Part 2 · Prior work")
two_col(s,
    "Interactive / failure-driven",
    [
        (0, "DAgger & interactive IL: expert corrects the policy at the states *it actually visits*"),
        (0, "Hard-example mining / failure replay: over-sample where the model errs"),
        (0, "Right instinct — target failures — but needs an expert in the loop at every step"),
    ],
    "The gap this work targets",
    [
        (0, "Validate selection on **closed-loop rollout success**, not held-out loss"),
        (0, "Test **cross-policy** — is a weak region universal or one policy's quirk?"),
        (0, "Run signals **head-to-head** under an identical recipe"),
        (0, "Be **honest**: power analysis + adversarial checks, report the negatives"),
    ],
    lcolor=MUTED, rcolor=TEAL)
_notes(prs.slides[-1], """
Family three: correct it. DAgger and interactive imitation learning put an expert in the loop to
correct the policy at the states it actually visits — which is exactly the right instinct, target the
failures — but it needs that expert online at every step, which is expensive and often infeasible.

So here's the gap I set out to fill, on the right. I wanted to (1) validate data selection on actual
closed-loop rollout success, not a proxy loss; (2) check whether a weak region is universal across
very different policies or just one model's quirk; (3) put the selection signals head-to-head under a
strictly identical training recipe so only the data differs; and (4) be ruthlessly honest — do the
power analysis, run adversarial checks, and report the negative results. That honesty is what made the
project actually informative.
""")

# ==========================================================================
# SECTION 3
# ==========================================================================
section_slide("Part 3", "My theory — and why it's unique",
              "Prescribe from the policy's own failure regions — then test it the hard way.")
_notes(prs.slides[-1], "Two and a half minutes: the thesis and the five things that make the setup unusual.")

# Slide: thesis
s = content_slide("The thesis",
                  kicker="Part 3 · Theory")
bullets(s, [
    (0, "==A trained policy's own failure regions are a prescription for what data to collect next.=="),
    (0, "Concretely: evaluate the policy, localize *where* and *how* it fails, then add demonstrations that **densify exactly those regions** — and expect more improvement-per-demo than random collection"),
    (0, "Testable sub-claims, built up in sequence:"),
    (1, "Failures are **structured** (localizable), not random noise"),
    (1, "The weak region is a **shared data gap**, not a hardware limit — so data *can* fix it"),
    (1, "**Concentrated** targeted data beats random; and beats a **principled** influence signal?"),
    (0, "The honest test: hold the training recipe fixed, change **only** the 200 selected demos, measure real rollout success"),
])
_notes(prs.slides[-1], """
My thesis in one sentence: a trained policy's own failure regions are a prescription for what to
collect next. Evaluate, localize where and how it fails, then add demonstrations that densify exactly
those regions — and you should get more improvement per demo than collecting at random.

That breaks into a chain of testable claims, which is how I'll structure the evidence. First, are
failures even structured — can you localize them? Second, is the weak region a fixable data gap or a
hardware limit? Third, does concentrated targeted data beat random — and can a fancy influence method
beat a dumb heuristic?

And the discipline that holds it together: fix the training recipe completely, change only the 200
demonstrations you select, and measure actual rollout success. Same recipe, only the data differs.
""")

# Slide: why unique
s = content_slide("Why this setup is unique",
                  kicker="Part 3 · Theory")
two_col(s,
    "Methodology",
    [
        (0, "Closed-loop: real simulator + real rollout eval (not held-out loss)"),
        (0, "Identical-recipe invariant: only the 200-demo arm changes"),
        (0, "Power analysis up front — knows what it can and can't detect"),
        (0, "Adversarial verification of every positive claim"),
    ],
    "Scope",
    [
        (0, "Cross-policy: pi0 (flow-matching) AND GR00T (diffusion)"),
        (0, "Head-to-head signals: P(fail) · coverage · gradient-influence · visual retrieval (CLIP/DINOv2)"),
        (0, "A clean sandbox (CIFAR) to isolate *why* methods work"),
        (0, "Reports the negative results as first-class findings"),
    ],
    lcolor=TEAL, rcolor=NAVY)
note_strip(s, "The payoff of this rigor: the most valuable result — why the sophisticated methods failed — only became visible because of the identical-recipe invariant and the sandbox.")
_notes(prs.slides[-1], """
What makes this different from a typical data-selection paper is on this slide.

On methodology: everything is closed-loop — a real RoboCasa simulator with real rollout evaluation,
not a held-out loss proxy. Every experiment obeys an identical-recipe invariant: same base weights,
same LoRA config, same steps, same everything — the *only* thing that changes is which 200 demos I add.
I did a power analysis up front, so I know what effect sizes I can even detect. And I adversarially
verified every positive claim rather than trusting the first happy number.

On scope: I test cross-policy — two architecturally very different policies, pi0 which is flow-matching
and GR00T which is diffusion. I run four selection signals head-to-head. And I built a clean CIFAR
sandbox to isolate *why* methods work when the robot experiments were too noisy to tell.

The payoff of all that rigor: the single most valuable finding — why the sophisticated methods lost —
only became visible *because* of the invariant and the sandbox. Rigor wasn't bureaucracy here; it was
the instrument.
""")

# ==========================================================================
# SECTION 4
# ==========================================================================
section_slide("Part 4", "Evidence — what matters, and what doesn't",
              "Five hypotheses, five experiments. Each says what mattered — and what surprisingly didn't.")
_notes(prs.slides[-1], "The core, eight minutes. Five hypothesis/experiment/result triples. I foreground both the wins and the null results.")

# Slide: testbed
s = content_slide("The testbed",
                  kicker="Part 4 · Setup")
two_col(s,
    "Environment & policies",
    [
        (0, "RoboCasa · task PickPlaceCounterToSink (rich object/pose variation)"),
        (0, "Students: pi0 (~55–58%) and GR00T (~56–66%) — public VLA checkpoints"),
        (0, "Pool: 9,885 MimicGen demos over 79 object categories"),
        (0, "Fine-tune: LoRA, 20k steps; base 400 demos + a 200-demo selected arm"),
    ],
    "The measurement discipline",
    [
        (0, "Eval = rollout success, n=300, fixed seed, scene-paired across arms"),
        (0, "Invariant: every arm is dataclasses-identical except its data_dirs"),
        (0, "No human demos needed — expert-teacher loop generates trainable data"),
        (0, "Report overall + targeted-10 + non-targeted success separately"),
    ],
    lcolor=NAVY, rcolor=TEAL)
_notes(prs.slides[-1], """
Quick setup so the numbers land. One RoboCasa task, PickPlaceCounterToSink — deliberately one task, so
I can debug it, but it has real variation in object type, pose, grasp, and placement. My 'students' are
pi0 and GR00T, off-the-shelf checkpoints around 55%. The candidate data pool is ~9,900 MimicGen demos
across 79 object categories.

Every experiment fine-tunes with LoRA for 20k steps on 400 base demos plus a 200-demo *selected* arm.
Evaluation is rollout success — n=300, fixed seed, scene-paired so arms see the same initial states.
The invariant, again: the arms are literally identical dataclasses except for which folder the 200 demos
come from. And I always split results three ways: overall, the targeted-10 categories, and the
non-targeted majority — because, spoiler, that split is where the whole story lives.
""")

# ---- H1 ----
s = content_slide("Hypothesis 1 — failures are structured, not random",
                  kicker="Part 4 · Experiment 1",
                  subtitle="Exp: pi0, n=150 rollouts + per-episode geometry + logistic predictor + embodiment test")
two_col(s,
    "What matters ✓",
    [
        (0, "One failure MODE dominates: ==the GRASP=="),
        (1, "76% of failures never touch the object (moves <1cm)"),
        (0, "Object HEIGHT is a real, moderate predictor"),
        (1, "5-fold CV AUC 0.63; tall (>11cm) 36% vs short 67%"),
        (0, "==Not an embodiment limit==: objects fit the gripper → skill/data gap, EPISTEMIC & fixable"),
    ],
    "What DOESN'T matter ✗",
    [
        (0, "Simple geometry is a WEAK targeting axis"),
        (1, "All geometric features together: R² ≈ 0.08"),
        (0, "Width, position, depth: coefficients ≈ 0"),
        (0, "Naive eyeballing at n=50 said 'tall fails 0%' — that was small-sample NOISE"),
        (1, "~90% of failure variance is object-instance-specific, not geometry"),
    ],
    lcolor=TEAL, rcolor=AMBER)
_notes(prs.slides[-1], """
Hypothesis 1: failures are structured — you can localize them. Experiment: run pi0 for 150 rollouts,
log per-episode object geometry, fit a logistic predictor, and run an embodiment test.

Left, what matters. There's a single dominant failure mode: the grasp. 76% of failures never even touch
the object — the policy fails at *initiating* the grasp, not at transport or placement. And object height
is a real predictor: cross-validated AUC 0.63, tall objects succeed 36% versus short objects 67% — a clean
30-point gap. Critically, I ran an embodiment test: is this a hardware limit? No — the objects fit the
gripper aperture; they *are* graspable. So this is an epistemic, data-addressable skill gap. Good news for
the whole premise.

But right — and this is the first 'what doesn't matter.' Simple geometry is a *weak* handle. All the
geometric features together explain only 8% of the variance. Width, position, depth — essentially zero.
And a cautionary tale: at n=50, eyeballing said 'tall objects fail 100% of the time.' That was pure
small-sample noise; it evaporated at n=150. Roughly 90% of the failure variance is specific to the
individual object instance, not to any tidy geometric feature. Lesson: localize the *mode* robustly, but
don't over-trust hand-picked geometry.
""")

# ---- H2 ----
s = content_slide("Hypothesis 2 — the weak region is universal",
                  kicker="Part 4 · Experiment 2",
                  subtitle="Exp: replicate the weak-region analysis on GR00T (a completely different architecture), n=100")
result_table(s, [
    ["Metric", "pi0  (flow-matching, n=150)", "GR00T  (diffusion, n=100)"],
    ["Overall success", "52.7%", "56.0%"],
    ["Dominant failure mode", "86% no-grasp", "80% no-grasp"],
    ["Geometry→success CV AUC", "0.628", "0.629"],
    ["Dominant predictor (coeff)", "height (−0.57)", "height (−0.78)"],
    ["Tall-object success", "~36%", "~32%"],
    ["Short-object success", "~67%", "~66%"],
], top=Inches(2.15), col_ratios=[1.5, 1.4, 1.4], height=Inches(3.4))
note_strip(s, "✓ Two very different architectures fail at the same rate, same mode, same objects → the tall-object grasp weakness is UNIVERSAL = a shared training-DATA gap, not a pi0 quirk.  ✗ What doesn't matter: policy disagreement — they agree too much to mine as a signal.", top=Inches(5.9))
_notes(prs.slides[-1], """
Hypothesis 2: is this weak region universal, or just a pi0 quirk? If it's universal across very
different models, that strongly implies a shared *data* gap rather than one model's idiosyncrasy.
Experiment: run the exact same weak-region analysis on GR00T — a completely different architecture,
diffusion with an Eagle VLM, versus pi0's flow-matching.

Look at this table. It's almost eerie. Same overall success, ~53 vs 56. Same dominant failure mode,
grasp, ~80–86%. Same geometry-to-success AUC, 0.628 vs 0.629. Same dominant predictor, height, with a
similar coefficient. Same tall-object success, ~32–36%. Same short-object success, ~66%.

Two architectures that share almost nothing internally fail in exactly the same way on exactly the same
objects. Combined with the embodiment test from before, that rules out 'it's the gripper' and rules out
'it's a pi0 bug.' It's a shared gap in the training data — both were trained on RoboCasa data that
under-serves tall-object grasps. That is precisely the kind of thing prescribed data should fix.

The one 'doesn't matter': I'd hoped to mine policy *disagreement* as an acquisition signal, but they
agree far too much — there's no disagreement to exploit here.
""")

# ---- H3 (reframed: overall-first) ----
s = image_slide(
    "Hypothesis 3 — read the OVERALL number",
    "Part 4 · Experiment 3",
    FIGDIR + "/fig_overall.png",
    subtitle="Only the 200-demo arm differs: core (P(fail) top-10) vs random vs coverage · n=300",
    bullets_items=[
        (0, "==Overall success is the metric we can trust== — full n=300, no sub-region cherry-picking"),
        (0, "**core 0.593 is the only arm above baseline** (0.580)"),
        (0, "Concentration matters: **coverage** (spread over top-25) is **worst**, 0.517 — thin spread → forgetting"),
        (0, "random / failure-influence / value all land ≤ baseline"),
        (0, "//But 'did it help the targeted region?' is a shakier question — next 3 slides//"),
    ],
    caption="The robust result is this overall ranking. Whether an arm 'helped the targeted region' turns out to depend on how you define that region — which is itself the problem.",
)
_notes(prs.slides[-1], """
The central applied experiment. Same recipe, only the 200-demo arm changes: core concentrates on the
top-10 failure categories, random is the control, coverage spreads the same budget over the top-25.

I want to lead with the metric we can actually trust — overall success, the full n=300, no sub-region
slicing. On that: core at 0.593 is the *only* arm that beats the no-fine-tune baseline at 0.580. And
concentration matters — coverage, spreading the *same* budget wider, is the *worst* arm at 0.517: thin
spread causes broad drift and forgetting. Everything else — random, influence, value — lands at or below
baseline.

Now, you'll notice I'm *not* leading with 'core lifted the targeted failure region.' I used to. But when
I stress-tested that claim it turned out to be fragile — and that fragility is itself one of the more
honest findings in this work. The next three slides show why the overall number, not a targeted slice, is
the claim I'll stand behind.
""")

# ---- H3b: the targeted region is fragile (both definitions) ----
s = image_slide(
    "…but the 'targeted region' is fragile",
    "Part 4 · Robustness",
    FIGDIR + "/fig_target_instability.png",
    subtitle="Four ways to score the same failure region — the verdict swings, then VANISHES when you select AND measure it cleanly",
    caption="Redefine which 10 categories are 'hard' → core's lift over baseline flips (+0.10 → −0.01). Re-sample them properly → core-vs-random collapses to a tie. Do BOTH right (balanced-select + balanced-measure, rightmost) → baseline 0.33 · core 0.35 · random 0.33 — all tied within noise. 'Targeting helped the failure region' doesn't survive clean measurement, so the OVERALL number carries the claim.",
)
_notes(prs.slides[-1], """
Here's the fragility. Four ways to score success on the *same failure region*, and the verdict swings —
then collapses.

Group one: the original targeted-10, sliced from the n=300 eval. Core beats baseline by about +0.10 — the
number I used to headline. Group two: keep the n=300 eval but re-derive the 10 hard categories from a
*balanced* Wilson-lower-bound ranking — five of the ten swap out — and that core-over-baseline lift flips
to −0.01. Gone. Group three: take the *same* old categories but measure them on a properly stratified eval,
tens of episodes per category instead of the natural handful. Now core and random are basically tied —
0.37 versus 0.35.

And group four — the one I built after being challenged on exactly this — does it *all* the right way:
select the hard categories from a balanced eval AND measure them with balanced per-category sampling. When
you do that, everything ties: baseline 0.33, core 0.35, random 0.33. Core edges the others by two points,
well inside noise. So "targeting helped the failure region" simply does not survive a clean measurement.
The verdict was an artifact of how the region was picked and sampled — while the overall number, from the
previous slide, never moved. That's why I show the targeted slice only for context and stand behind the
overall number.
""")

# ---- H3c: per-category heatmap ----
s = image_slide(
    "Per-category texture: targeting is uneven",
    "Part 4 · Per-category",
    FIGDIR + "/fig_percat_heatmap.png",
    subtitle="Stratified per-category success (per_cat ≈ 16–39) — and it misses the genuinely-hardest",
    caption="core helps MEDIUM cats most (canned_food +0.32, tupperware +0.23, spray +0.21) but barely moves the genuinely-hardest (juice −0.01; cheese_grater still 0.19). random helps too (tupperware +0.42, beats core). whiten ≈ baseline everywhere → its ranking win did not translate.",
)
_notes(prs.slides[-1], """
Now the per-category texture — the heatmap you asked to see. Rows are the categories, hardest at the top;
the right panel is each arm's change versus baseline, green for helped, red for hurt.

Two things jump out. First, core's help is *uneven* and, tellingly, it barely touches the genuinely
hardest categories: juice is −0.01, cheese_grater goes from 0.04 to 0.19 — still almost always failing.
Where core helps is the *medium* categories — canned_food plus 0.32, tupperware plus 0.23, spray plus 0.21
— categories that weren't that hard to begin with. Second, random helps in several places too — it even
beats core on tupperware, plus 0.42 — which is another way of seeing that core-over-random is not a clean
win. And whiten, the far column, is essentially baseline everywhere: its ranking-AUC win did not translate
into per-category rollout success.
""")

# ---- H3d: the real weak axis (height) ----
s = image_slide(
    "The real weak axis is height",
    "Part 4 · The deeper cut",
    FIGDIR + "/fig_height.png",
    subtitle="Every intervention arm REGRESSED the tall objects — the actual weak region",
    bullets_items=[
        (0, "H1/H2: **tall objects are the actual weak region** (the grasp fails)"),
        (0, "Baseline tall = 0.52; **every intervention falls BELOW** (core 0.43, random 0.35, all arms)"),
        (0, "The category-targeted data helped short/medium — **not** the height-defined weak spot"),
        (0, "==Category is a poor proxy for the real weak region== (recall H1: geometry R²≈0.08)"),
        (0, "Consistent across all 5 arms → not noise"),
    ],
    caption="The prescription improved where it was needed least and HURT the genuinely-hard tall grasps — the strongest reason to trust the overall number over any category slice.",
)
_notes(prs.slides[-1], """
And here's the deeper cut, tying back to Hypotheses 1 and 2. The real weak axis isn't category — it's
object height; tall objects are where the grasp fails. So slice the same evals by height tertile.

Baseline is the dark line. Look at the tall end on the right: baseline sits at 0.52, and *every single
intervention arm falls below it* — core at 0.43, random at 0.35, all of them. The extra 200 demos, whatever
the selection, improved the short and medium objects and *regressed* the tall ones — the actual weak region.
And it's consistent across all five arms, so it's not noise.

That's the sharpest statement of the whole section: category-based prescription helped where it was needed
least and hurt where it was needed most. It's the strongest reason to report the overall number and treat
any single category slice with suspicion — which is exactly the fragility Hypothesis 1 predicted, where
geometry explained only eight percent of the variance.
""")

# ---- H4 ----
s = content_slide("Hypothesis 4 — a principled influence signal should beat the heuristic",
                  kicker="Part 4 · Experiment 4",
                  subtitle="Exp: score every pool demo by gradient influence (LESS/TracIn); two variants; fine-tune the top-200; eval n=300")
result_table(s, [
    ["Arm", "OVERALL", "Targeted (old)", "Targeted (bal.)", "Selection"],
    ["core — trivial P(fail) heuristic", "0.593", "0.432", "0.432", "100% targeted"],
    ["baseline", "0.580", "0.333", "0.439", "—"],
    ["random", "0.563", "0.282", "0.235", "~7% targeted"],
    ["failure-influence — contrast g_val", "0.553", "0.421", "0.302", "23% targeted"],
    ["coverage", "0.517", "0.429", "0.289", "spread top-25"],
    ["value-influence — plain g_val", "0.503", "0.361", "0.359", "6% targeted"],
], top=Inches(2.1), col_ratios=[2.0, 1.0, 1.15, 1.15, 1.15],
   highlight_rows={1: (RGBColor(0xE3, 0xF1, 0xF0), NAVY),
                   4: (RGBColor(0xF7, 0xE7, 0xDE), AMBER),
                   6: (RGBColor(0xF3, 0xD9, 0xCE), AMBER)}, height=Inches(3.0))
note_strip(s, "Read the OVERALL column — it's the robust one. ✗ BOTH gradient-influence methods LOSE to the trivial heuristic AND to baseline. The two targeted columns disagree (old vs balanced), so the targeted slice is not the claim.", top=Inches(5.4))
note_strip(s, "value-influence picked data ~disjoint from failures (6% targeted) → evidence against 'failure ⇒ missing data'. Only core clears baseline overall — by doing least damage, not by being smart.", top=Inches(6.55))
_notes(prs.slides[-1], """
Hypothesis 4 is where I expected to show off. The 'core' arm is a dumb heuristic — just the top-10
categories by failure rate. Surely a *principled* method — gradient influence, LESS/TracIn — does better?
The idea: score every demo in the pool by how much its gradient aligns with a target gradient, take the
top-200. I built two variants. Failure-influence aims at the failure region using a contrast direction.
Value-influence aims at improving overall balanced success.

The result, honestly, humbled me. Both influence methods *lose* — not just to core, but to baseline.
Failure-influence lands at 0.553, value-influence at 0.503, the worst of everything. Read the table:
core 0.593, baseline 0.580, random 0.563, failure-influence 0.553, coverage 0.517, value-influence 0.503.
The trivial heuristic beat both sophisticated selectors.

Two diagnostic facts. Failure-influence actually worked *where it aimed* — 0.421 on the targeted region,
essentially matching core — but it regressed the non-targeted majority and that ate the gains. And
value-influence, when asked 'what data most improves overall success,' picked demos almost entirely
*outside* the failure region — only 6% targeted, less than random. So 'the valuable data' and 'the
failure data' were nearly disjoint sets.

And here's the unifying observation, the bottom line: the non-targeted column predicts the entire ranking.
Every arm that perturbs the shared grasp behavior forgets, and loses overall. Only core leaves it intact.
The heuristic won by *doing less damage*, not by being smarter. That demanded an explanation — Hypothesis 5.
""")

# ---- H4b: visual-driven retrieval (CLIP + DINOv2) ----
s = image_slide(
    "Visual-driven retrieval (2.3 CLIP · 2.4 DINOv2)",
    "Part 4 · Experiment 4b",
    FIGDIR + "/fig_selection_targeting.png",
    subtitle="Select demos whose scene looks like where pi0 fails — CLIP mean-top-k · DINOv2 greedy facility-location",
    bullets_items=[
        (0, "**2.3 CLIP retrieval**: mean top-k cosine of each pool frame to pi0's failure scenes → top-200"),
        (0, "**2.4 DINOv2 greedy coverage**: facility-location, max marginal coverage of the failure set"),
        (0, "==Both select ≈ random on the failure region== (CLIP 10%, DINOv2 9% vs random 6%)"),
        (0, "Scene embeddings ≈ the shared kitchen **common-mode** → 'looks like a failure scene' ≠ 'the hard object'"),
        (0, "**Rollout-confirmed**: failretr (DINOv2-contrast) → **0.297** targeted-10 (< random 0.351 · < core 0.371)"),
    ],
    caption="Same common-mode / encoding failure as the gradient methods — in VISUAL feature space. Even the one visual arm carried to rollout (failretr) lands below random.",
)
_notes(prs.slides[-1], """
One more family I want to be transparent about, because it's the natural thing to try and you asked to
see it: visual-driven retrieval. The intuition is compelling — instead of scoring demos by gradients,
just find the pool demos whose *scene looks like* the scenes pi0 actually fails in. I built two. 2.3 uses
CLIP: embed pi0's failure scenes, embed every pool demo's first frame, score each demo by its mean top-k
cosine similarity to the failure set, take the top 200. 2.4 uses DINOv2 with a greedy facility-location
objective — add, one at a time, the demo that most increases coverage of the failure distribution.

Here's the selection-level result, and it's the same story from a new direction. This chart is 'what
fraction of each method's 200 picks actually fall in the failure region.' core is 100% by construction.
But CLIP lands at 10%, DINOv2-greedy at 9%, the DINOv2-contrast variant at 8% — all essentially random,
which is 6%. Visual similarity to failure scenes does *not* recover the hard objects. Why? Because a
whole-scene embedding — the counter, the layout, the robot arm — is dominated by the shared kitchen
common-mode; 'looks like a failure scene' is mostly 'looks like a kitchen,' not 'is the specific tall
object that fails.' It's exactly the common-mode problem from the gradient methods, now in visual feature
space. To be clear on status: the CLIP and DINOv2-greedy variants were selection-tested only. But the
third visual variant — DINOv2 fail-minus-success contrast, "failretr" — I did carry all the way to a
fine-tune and rollout eval, and it confirms the prediction: 0.297 on the targeted region, below random
at 0.351 and core at 0.371, barely above baseline. So visual retrieval doesn't just fail to concentrate
on failures at the selection stage — when you actually train on it, it loses too. That sets up the
unifying principle on the next slide.
""")

# ---- H5 ----
s = content_slide("Hypothesis 5 — influence works iff the gradient encodes the target",
                  kicker="Part 4 · Experiment 5 · the synthesis",
                  subtitle="Exp: CIFAR sandbox (clean ground truth) + robocasa gradient diagnostics + whitening")
result_table(s, [
    ["Setting", "Does the loss gradient encode the target?", "Best SVD-mode AUC", "Influence AUC"],
    ["CIFAR — select by class (cross-entropy on class)", "YES — target IS the loss", "0.66", "0.96  ✓"],
    ["RoboCasa — select by object (flow-loss on ACTIONS)", "NO — object is a nuisance to action loss", "0.56", "0.60  ✗"],
    ["CIFAR — select by brightness (a nuisance var)", "NO — nuisance to classification loss", "0.53", "0.50  ✗"],
], top=Inches(2.15), col_ratios=[2.0, 2.0, 1.0, 1.0],
   highlight_rows={1: (RGBColor(0xE3, 0xF1, 0xF0), NAVY),
                   2: (RGBColor(0xF7, 0xE7, 0xDE), AMBER),
                   3: (RGBColor(0xF7, 0xE7, 0xDE), AMBER)}, height=Inches(2.6))
bullets(s, [
    (0, "==The principle==: gradient-influence selection works **iff the loss gradient encodes the target distinction** — predictable in advance from the best single-SVD-mode AUC of the gradient cloud"),
    (0, "RoboCasa fails because **object identity is a nuisance** — to the action *loss* (gradients) AND to whole-scene *embeddings* (CLIP/DINOv2 select ≈random): a common-mode nuisance in **both** feature spaces"),
    (0, "==The one method win — whitening==: project out the top shared 'common-mode' modes → robocasa ranking AUC **0.605 → 0.677**, targeted purity 17% → 36% (must bound k)"),
    (0, "**…but rollout ≠ ranking**: whitening fine-tuned + rollout-evaluated → **0.268 ≈ baseline 0.262 ≪ core 0.371** → the ranking win did **not** translate (exactly as the honest hedge predicted)"),
], top=Inches(4.7), base_size=15)
_notes(prs.slides[-1], """
Hypothesis 5 is the synthesis, and for me the real contribution. Why did influence work beautifully in
the LLM/vision literature but flop on my robot task? To isolate it I built a clean CIFAR sandbox with a
ground-truth 'failure region' — 20 rare classes — where I control everything.

The table is the whole story. When I select CIFAR images by *class*, and the loss is cross-entropy *on
that class*, influence gets AUC 0.96 — it works brilliantly. When I select CIFAR images by *brightness* —
a nuisance variable the classification loss doesn't care about — influence collapses to 0.50, chance.
And robocasa sits in the middle at 0.60, and now we know why: I'm selecting by *object identity*, but the
loss is flow-matching on *actions*, and object identity is a *nuisance* to an action-prediction loss —
structurally the same situation as brightness for a classifier.

So the principle: gradient-influence selection works if and only if the loss gradient actually encodes the
distinction you're targeting. And you can predict it cheaply in advance — the best single SVD-mode AUC of
the gradient cloud tells you the ceiling before you spend a GPU-hour. 0.66 for CIFAR-class, 0.56 for
robocasa-object, 0.53 for the nuisance.

There was one genuine method win: whitening. Project out the top shared 'common-mode' gradient directions —
the generic pick-place component every demo shares — and the robocasa ranking improves from 0.605 to 0.677,
with targeted purity doubling from 17 to 36%. But you must bound how many modes you remove; over-whiten and
you delete the signal too.

And I can now close the loop honestly: I fine-tuned the whitening selection and rollout-evaluated it. It
lands at 0.268 — essentially the baseline's 0.262, and well below core's 0.371. So the ranking-AUC win did
*not* translate into rollout success, exactly as I hedged it might not. A better *ranking* of demos is not
the same as a better *policy* — which is the honest boundary on this whole line of work.
""")

# ---- what matters / doesn't matter summary ----
s = content_slide("Evidence, distilled: what matters vs. what doesn't",
                  kicker="Part 4 · Synthesis")
two_col(s,
    "What MATTERS / is TRUSTWORTHY ✓",
    [
        (0, "The **OVERALL** number — one metric, no sub-region cherry-picking"),
        (0, "Localizing the failure MODE (grasp) — robust & cross-policy"),
        (0, "A data-addressable gap (not a hardware limit)"),
        (0, "CONCENTRATION over spread (coverage was worst)"),
        (0, "Preserving the shared skill (do-no-harm on the majority)"),
        (0, "Matching the selection SIGNAL to what the loss encodes"),
    ],
    "What's MISLEADING / backfires ✗",
    [
        (0, "The 'targeted-region' metric — 2 category defs disagree (5/10)"),
        (0, "Category as a weak-region proxy (geometry R²≈0.08)"),
        (0, "Targeting by category HURT the real (tall) weak region"),
        (0, "Influence AND visual retrieval when the target isn't encoded (gradient or scene)"),
        (0, "Whitening's ranking win (0.677) — did NOT translate to rollout"),
        (0, "'Value' selection that drifts off the failure region"),
    ],
    lcolor=TEAL, rcolor=AMBER)
_notes(prs.slides[-1], """
Let me distill the section — the 'what matters / what doesn't' the title promised, now sharpened by the
robustness checks.

What's trustworthy: the *overall* number, first and foremost — one metric, no sub-region games. Then the
failure *mode*, the grasp, which is robust and cross-policy; the fact that it's a data gap, not hardware;
concentration over spread; preserving the shared skill; and matching the selection signal to what the loss
encodes.

What's misleading or backfires: the targeted-region metric itself — two reasonable category definitions
disagree five out of ten, so it can't carry the claim. Category is a weak proxy for the real weak region —
geometry explained eight percent of the variance — and targeting by category actually *hurt* the real,
height-defined weak region. Sophisticated influence when the gradient can't see the target. Whitening's
ranking win that didn't translate to rollout. And value-style selection that drifts off the failure region.

The through-line: trust aggregate, distribution-preserving signals; distrust any narrow sub-region metric —
including the one I started this project believing in.
""")

# ---- real-robot corroboration (bridge Part 4 -> Part 5) ----
s = content_slide("A real-robot echo — Unitree G1 chemistry pouring",
                  kicker="Part 4 · Corroboration",
                  subtitle="The pouring task the sim study stood in for — same headline on a real humanoid, via a different selection axis")
two_col(s,
    "The setup",
    [
        (0, "Unitree G1 humanoid + Dex3 hands; pour a tube into a beaker"),
        (0, "Teleop → LeRobot → diffusion policy; 28-D state/action, wrist cam @30fps"),
        (0, "Tiny, expensive data: 11–26 teleop episodes = minutes of motion"),
        (0, "Same H100 stack; policies: diffusion / ACT / pi0 / GR00T"),
    ],
    "Prescribed selection (SREE clustering)",
    [
        (0, "Full 26-episode pool → diffusion MODE-COLLAPSES, freezes mid-episode"),
        (0, "PCA + pairwise-MSE clustering → keep the consistent 8-episode cluster"),
        (0, "−65% frames (18,004 → 6,272), outliers dropped"),
        (0, "==Result: clean single-mode convergence, no freeze — quality > quantity=="),
    ],
    lcolor=NAVY, rcolor=TEAL)
note_strip(s, "Honest caveat: this evidence is qualitative + training-loss (mode collapse fixed), NOT yet a rollout-success comparison — exactly the rigorous eval Part 5 calls for. But the headline repeats the sim result: which demonstrations you train on beats how many.")
_notes(prs.slides[-1], """
Before the path forward, let me close the loop back to that real robot — because the same lesson shows
up on hardware. This is a Unitree G1 humanoid with dexterous hands, learning the chemistry-pouring task
I mentioned at the start: pick up a tube, pour it into a beaker. Data is collected by human teleoperation
in a headset — each episode is minutes long, so the whole dataset is tiny: 11 to 26 episodes.

Here's the prescribed-data moment. When we trained a diffusion policy on the *full* 26-episode pool, it
mode-collapsed — the policy froze mid-episode, because the demos came from two inconsistent motion styles
plus a couple of outliers. So instead of dumping all the data in, we clustered the demonstrations — PCA
plus pairwise trajectory MSE — and kept only the tight, consistent 8-episode cluster. That's a 65% cut in
frames. And the policy converged cleanly to a single mode, no more freezing.

Now, the honest caveat, and it's important given everything I said about rigor: this evidence is
qualitative plus training-loss — the mode collapse went away — not yet a controlled rollout-success
comparison. That's exactly the rigorous evaluation my path-forward calls for. But notice the headline is
the *same* as the simulator study, reached through a different selection axis — demonstration consistency
instead of failure targeting: which demonstrations you train on matters more than how many. Prescribed
data isn't a sim artifact; it's already earning its keep on a real humanoid.
""")

# ==========================================================================
# SECTION 5
# ==========================================================================
section_slide("Part 5", "The path forward",
              "Point the gradient at what the loss encodes; power up the eval; close the real-robot loop.")
_notes(prs.slides[-1], "Two and a half minutes: where this goes next, ordered by leverage.")

# Slide: make influence work
s = content_slide("Making prescribed selection actually work",
                  kicker="Part 5 · Next")
bullets(s, [
    (0, "==Re-aim the gradient at what the action loss encodes== — not object identity:"),
    (1, "grasp-success / motion-phase / per-scene targets; bar to clear = best-mode AUC ≳ 0.7"),
    (0, "**RC-LESS** — the stacked selector: retention-constrained + coverage-aware (facility-location) + Adam-preconditioned + whitened, with a per-category floor"),
    (1, "By construction it never does worse than core; it can only win via within-category structure core is too coarse to see"),
    (0, "**Gate on success, not loss**: a cheap 2k-step mini-fine-tune that rejects MSE-good / success-bad selections before any full run"),
    (0, "Honest ceiling: this may be a **data-coverage** limit, not an algorithm limit — when value-data lies outside the failure region, only *collecting different demos* clears core"),
], base_size=17)
_notes(prs.slides[-1], """
So where does this go? Ordered by leverage.

Number one, the direct implication of Hypothesis 5: re-aim the gradient at something the action loss
actually encodes — grasp-success, motion-phase, per-scene identity — instead of object category. I even
have a quantitative bar: get the best-mode AUC above about 0.7 before bothering to run the selection.

Two, the concrete next method — I call it RC-LESS. It stacks the fixes that address orthogonal defects:
retention-constrained so it doesn't forget, coverage-aware via facility-location instead of naive top-k,
Adam-preconditioned to respect gradient magnitude, whitened to kill the common mode — and a per-category
floor that guarantees it *never does worse than core*, while giving it a shot at winning through
within-category structure that core, being category-coarse, simply cannot see.

Three, and this is the deepest lesson from the failures: gate on *success, not loss*. A cheap 2k-step
mini-fine-tune that rejects the selections that look good on MSE but tank real rollout success — the exact
trap value-influence fell into.

And I want to be honest about the ceiling: my analysis suggests this might be a *data-coverage* limit, not
an algorithm limit. When the most valuable data genuinely lies outside the failure region, no clever
re-selection of the existing pool beats core — only going out and *collecting different demonstrations* does.
Which loops right back to the point of prescribed data.
""")

# Slide: rigor + real robot
s = content_slide("From simulator to real robots",
                  kicker="Part 5 · Next")
two_col(s,
    "Power up the evidence",
    [
        (0, "Paired-seed design + ~16k rollouts to resolve 1–4 pt gaps (n=300 can't)"),
        (0, "Multiple tasks & seeds — beyond single-task, single-run"),
        (0, "Fix the rollout-eval hang (per-rollout timeout) for clean large-n runs"),
    ],
    "Close the real-robot loop",
    [
        (0, "The G1 pouring pipeline is already built (teleop → LeRobot → train → eval)"),
        (0, "Next: a proper rollout-success eval on the robot (not just training-loss / mode-collapse)"),
        (0, "Generate prescribed data on demand via the expert-teacher loop (pi0/GR00T → trainable data)"),
        (0, "Prescribe → collect exactly that region → retrain → re-evaluate on the G1"),
    ],
    lcolor=NAVY, rcolor=TEAL)
_notes(prs.slides[-1], """
Two more threads. On evidence: n=300 rollouts simply cannot resolve the 1-to-4-point gaps that separate my
arms — the power analysis says so. The fix is a paired-seed design and something like 16,000 rollouts,
across multiple tasks and seeds, not a single task and a single run. And a practical engineering item —
fix the rollout-eval hang with a per-rollout timeout so I can actually run large-n evals unattended.

On impact: the G1 pouring pipeline that motivated all this is already built end-to-end — teleop, LeRobot
conversion, training, and eval — and the clustering result you just saw is the first prescribed-data win on
it. The missing rigorous piece is the same one as in sim: a proper rollout-success eval on the robot, not
just training-loss. And a nice enabler I validated along the way — an expert-teacher loop where pi0 or
GR00T generate trainable demonstration data on demand, with no human teleoperation needed. That means once
I can prescribe a region, I can actually *fill* it. Prescribe, collect exactly that region, retrain,
re-evaluate on the G1 — the loop closes.
""")

# ==========================================================================
# CLOSE
# ==========================================================================
s = content_slide("Takeaways",
                  kicker="In closing")
bullets(s, [
    (0, "**1 — Measure prescription by OVERALL success, not a targeted slice.** The 'targeted region' is fragile (two category definitions disagree 5/10), and category-targeting even *hurt* the real height-defined weak region. On the trustworthy overall metric, concentrated failure-data (core) is the only arm above baseline — a small, honest win."),
    (0, "**2 — Selection only works when the gradient encodes the target.** A trivial P(fail) heuristic beat sophisticated LESS/TracIn influence, and whitening's ranking win didn't translate to rollout — because object identity is a nuisance to an action-prediction loss. Predict it in advance (best-mode AUC)."),
    (0, "**3 — The contribution is a validated diagnostic + a discipline of honesty** — not a SOTA selector: report aggregate metrics, distrust narrow sub-regions, and check the signal is even encodable."),
    (0, "==Don't just add more data — prescribe the right data, measure it honestly, and first check your training signal can even see it.=="),
], base_size=16)
note_strip(s, "Thank you — questions welcome.  ·  Testbed: RoboCasa · pi0 / GR00T · identical-recipe LoRA fine-tunes, rollout-evaluated.", top=Inches(6.7))
_notes(prs.slides[-1], """
To close, three takeaways in the sequence I promised.

One, the applied result, stated carefully: measure prescription by *overall* success, not a targeted
slice — because the targeted region turned out fragile, and category-targeting even hurt the real,
height-defined weak region. On that trustworthy overall metric, concentrated failure-data is the only arm
that clears baseline. It's a small, honest win — and the honesty is the point.

Two, the deeper and more transferable lesson: data selection only works when your training gradient
actually encodes the target. A trivial failure-rate heuristic beat sophisticated influence functions,
because on this task object identity is a nuisance to the action loss. And you don't have to find this out
the hard way — the best single-mode AUC of the gradient cloud predicts it before you train.

Three, framing the contribution honestly: this is a validated diagnostic and principle, not a new
state-of-the-art selector — plus one genuine method win in whitening, and an honest ceiling: this may be a
data-coverage limit, where the only way past the heuristic is to collect genuinely different data.

The through-line for anyone doing data-efficient learning: don't just add more data. Prescribe the right
data — and first check that your training signal can even see it. Thank you. I'm happy to take questions.
""")

# ==========================================================================
# APPENDIX
# ==========================================================================
section_slide("Appendix", "Backup — method & statistics",
              "Details for Q&A: full recipe, influence method, power analysis.")

s = content_slide("Appendix — influence method & why it's hard",
                  kicker="Backup")
bullets(s, [
    (0, "Score(z) = cos( ∇_LoRA flow-loss(z), g_val ), top-200; streamed (no projection) over ~50M LoRA-adapter dims"),
    (0, "K=8 frames/demo; loss masked to the 12 real action dims (pi0 pads 12→32 with noise targets)"),
    (0, "**Gradient target matters**: last-layer readout is motion-dominated & object-blind (AUC 0.44) → must use LoRA adapters"),
    (0, "**Reference checkpoint matters**: must be selection-neutral AND in-regime → a base-only LoRA warmup, not the trained arms, not cold pretrained pi0"),
    (0, "**g_val direction**: contrast (mean_hard−mean_ref) for failure; plain mean for value — plain collapses onto the common 'generic pick-place' mode (smoke AUC 0.35)"),
    (0, "Four coupled defects: common-mode collapse · magnitude-blindness · loss≠success · no retention/diversity term"),
], base_size=15)

# ---- Appendix: RC-LESS ----
s = content_slide("Appendix — RC-LESS: the full stack ties (or loses to) core",
                  kicker="Backup",
                  subtitle="The roadmap's 'single best next method' — retention-constrained, coverage-aware, sketch-based influence")
bullets(s, [
    (0, "**Two stages**: (1) sketch each demo's ==LoRA-adapter flow-matching gradient== (sparse Johnson–Lindenstrauss projection, K=8 frames, neutral warmup ref); (2) select from the sketches"),
    (0, "**Cluster** D_val into m=14 gradient modes → **coverage** score = top-3 cosine to modes (facility-location: help *some* mode, not the blurry average)"),
    (0, "**Retention penalty**: demote demos too parallel to g_R = the shared-grasp / forgetting direction (≈0.78 ∥ the task gradient)"),
    (0, "**Per-category floor + cap** on the targeted-10 → reimplements core's depth-on-holes ⇒ RC-LESS ≥ core *by construction*"),
], top=Inches(1.95), base_size=14)
result_table(s, [
    ["variant", "key knobs", "targeted-frac", "rollout n=300"],
    ["RC-LESS v1", "λ=1 · contrast-center ON", "0.60", "0.477  (worst)"],
    ["RC-LESS v2", "λ=0 · no-center · floor=20", "0.96", "0.590"],
    ["core (trivial heuristic)", "top-10 P(fail)", "1.00", "0.593"],
], top=Inches(4.35), col_ratios=[1.4, 1.8, 1.1, 1.2], height=Inches(1.75),
   highlight_rows={1: (RGBColor(0xF3, 0xD9, 0xCE), AMBER),
                   3: (RGBColor(0xE3, 0xF1, 0xF0), NAVY)})
note_strip(s, "v1's centering + retention penalty steer AWAY from the task gradient (g_R ≈ 0.78 ∥ it) → below baseline. v2 turns them off and floor=20×10 fills the whole 200-demo budget → RC-LESS DEGENERATES to core (0.590 ≈ 0.593). The best-engineered selector never beats the trivial heuristic ⇒ a data-coverage ceiling, not an algorithm one.", top=Inches(6.5))

s = content_slide("Appendix — power, honesty, and statistics",
                  kicker="Backup")
bullets(s, [
    (0, "Power analysis up front: detecting Δ=5% needs ~50–125 seeds; Δ=15% needs ~5–9 → n=300 resolves only large effects"),
    (0, "At n=300 paired McNemar: only **value < core is significant** (p≈0.02); most overall gaps (1–4 pts) are within noise"),
    (0, "Defensible claims: core is (weakly) best & the only arm beating baseline; coverage & value significantly regress the majority; influence did not beat the heuristic"),
    (0, "The 'height fails 0%' n=50 claim was refuted at n=150 — a worked example of small-sample overconfidence"),
    (0, "CIFAR sandbox gives clean ground truth (influence AUC 0.96–1.0) → proves the machinery is sound; robocasa's weakness is the *target encoding*, not a bug"),
    (0, "Every positive result adversarially re-derived; negative results reported as first-class findings"),
], base_size=15)

# ---------- save ----------
out = "/data/xinyua11/robocasa/talk/prescribed_data_talk.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
